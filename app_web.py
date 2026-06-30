import streamlit as st
import google.generativeai as genai
from pypdf import PdfReader
import json
import zipfile
import pandas as pd
from docx import Document
from pptx import Presentation
import io
import time # [MỚI] Thêm thư viện thời gian

# Cấu hình trang
st.set_page_config(page_title="Hệ Thống Hỗ Trợ Học Tập", layout="wide", page_icon="📚")

# ======================================================
# 1. CẤU HÌNH API & MODEL
# ======================================================
try:
    if "GOOGLE_API_KEY" in st.secrets:
        api_key = st.secrets["GOOGLE_API_KEY"]
    else:
        api_key = st.sidebar.text_input("Nhập Google API Key:", type="password")

    if api_key:
        # Danh sách model tối ưu nhất hiện tại
        model_options = [
            "gemini-3.5-flash",       # [KHUYÊN DÙNG] Nhanh, thông minh, cân bằng nhất
            "gemini-3.5-flash-lite",  # Siêu tiết kiệm limit
        ]
        
        selected_model = st.sidebar.selectbox(
            "🤖 Chọn Mô hình (Model):", 
            model_options,
            index=0
        )
        
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel(selected_model)
        
        st.sidebar.success(f"Đang dùng: {selected_model}")
        st.sidebar.info("💡 Mẹo: Nếu gặp lỗi 'Hết lượt', hãy đổi sang model khác trong danh sách.")
        
except Exception as e:
    st.error(f"Lỗi cấu hình: {e}")

# ======================================================
# 2. HÀM XỬ LÝ API AN TOÀN (QUAN TRỌNG)
# ======================================================
def goi_gemini_an_toan(prompt_input):
    """
    Hàm gọi AI có chức năng tự động bắt lỗi hết lượt (Rate Limit).
    Trả về: response object nếu thành công, hoặc None nếu thất bại.
    """
    try:
        if 'model' not in globals():
            st.error("Chưa cấu hình API Key!")
            return None
            
        response = model.generate_content(prompt_input)
        return response
        
    except Exception as e:
        error_msg = str(e)
        # Kiểm tra lỗi hết lượt (429 hoặc ResourceExhausted)
        if "429" in error_msg or "ResourceExhausted" in error_msg:
            st.toast("🚨 Hết lượt rồi! Đổi model đi bạn ơi!", icon="😫")
            st.error("🚨 **HẾT LƯỢT (QUOTA EXCEEDED)!**")
            st.warning(f"👉 Model **{selected_model}** đã quá tải. Vui lòng chọn model khác (ví dụ: gemini-2.5-flash-lite) ở thanh bên trái.")
        else:
            st.error(f"⚠️ Lỗi hệ thống: {e}")
        return None

# ======================================================
# 3. CÁC HÀM ĐỌC FILE
# ======================================================
def doc_pdf(file_bytes):
    try:
        reader = PdfReader(file_bytes)
        text = ""
        for page in reader.pages:
            t = page.extract_text()
            if t: text += t + "\n"
        return text
    except: return ""

def doc_word(file_bytes):
    try:
        doc = Document(file_bytes)
        text = []
        for para in doc.paragraphs:
            text.append(para.text)
        return "\n".join(text)
    except: return ""

def doc_pptx(file_bytes):
    try:
        prs = Presentation(file_bytes)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except: return ""

def doc_excel(file_bytes):
    try:
        df = pd.read_excel(file_bytes)
        return df.to_string()
    except: return ""

def xu_ly_file_upload(file_obj, ten_file):
    ten_file = ten_file.lower()
    if ten_file.endswith('.pdf'): return doc_pdf(file_obj)
    elif ten_file.endswith('.docx'): return doc_word(file_obj)
    elif ten_file.endswith('.pptx'): return doc_pptx(file_obj)
    elif ten_file.endswith('.xlsx') or ten_file.endswith('.xls'): return doc_excel(file_obj)
    return ""

# --- HÀM PHỤ TRỢ ---
def lay_json(text):
    text = text.replace("```json", "").replace("```", "").strip()
    s = text.find("[")
    e = text.rfind("]") + 1
    return text[s:e] if s != -1 and e != -1 else text

def lay_dot_code(text):
    text = text.replace("```dot", "").replace("```graphviz", "").replace("```", "").strip()
    s = text.find("digraph")
    if s != -1: return text[s:]
    return text

# ======================================================
# 4. GIAO DIỆN CHÍNH
# ======================================================
st.title("📚 Hệ Thống Học Tập Tích Hợp Gemini")

with st.sidebar:
    st.header("📂 Nạp tài liệu")
    st.caption("Hỗ trợ: PDF, Word, Excel, PowerPoint và ZIP")
    
    uploaded_files = st.file_uploader("Tải file lên:", type=['pdf', 'docx', 'pptx', 'xlsx', 'zip'], accept_multiple_files=True)
    
    if uploaded_files:
        if st.button("🔄 Xử lý tài liệu", use_container_width=True):
            with st.spinner("Đang đọc và phân tích..."):
                noi_dung_tong = ""
                ds_ten = []
                bar = st.progress(0)
                total_files = len(uploaded_files)
                
                for i, file in enumerate(uploaded_files):
                    if file.name.lower().endswith('.zip'):
                        try:
                            with zipfile.ZipFile(file) as z:
                                for sub_file in z.namelist():
                                    if not sub_file.startswith('__') and '.' in sub_file:
                                        with z.open(sub_file) as f_data:
                                            txt = xu_ly_file_upload(io.BytesIO(f_data.read()), sub_file)
                                            if txt:
                                                noi_dung_tong += f"\n--- FILE ZIP/{sub_file} ---\n{txt}\n"
                                                ds_ten.append(f"📦 {sub_file}")
                        except: st.error(f"Lỗi zip {file.name}")
                    else:
                        txt = xu_ly_file_upload(file, file.name)
                        if txt:
                            noi_dung_tong += f"\n--- FILE: {file.name} ---\n{txt}\n"
                            ds_ten.append(file.name)
                    
                    bar.progress((i + 1) / total_files)
                
                bar.empty()
                if ds_ten:
                    st.session_state['noi_dung'] = noi_dung_tong
                    st.session_state['ds_file'] = ds_ten
                    st.success(f"✅ Đã xử lý {len(ds_ten)} file!")
                else:
                    st.warning("Không tìm thấy nội dung.")

    if 'ds_file' in st.session_state:
        st.write("---")
        st.caption("Danh sách file:")
        for f in st.session_state['ds_file']:
            st.code(f, language="text")

# --- PHẦN TAB CHỨC NĂNG ---
if 'noi_dung' in st.session_state:
    t1, t2, t3, t4 = st.tabs(["💬 Chat", "📝 Trắc Nghiệm", "🗂️ Flashcards", "🧠 Sơ Đồ Tư Duy"])

    # -----------------------------------------------------
    # TAB 1: CHAT
    # -----------------------------------------------------
    with t1:
        if "msg" not in st.session_state: st.session_state.msg = []
        
        # Hiển thị lịch sử
        for m in st.session_state.msg: 
            with st.chat_message(m["role"]): st.markdown(m["content"])
        
        # [MỚI] Ghi chú Admin/Dev cho User
        st.warning("⚠️ **Lưu ý:** AI chỉ trả lời dựa trên tài liệu bạn đã nạp. Kiểm tra kỹ thông tin quan trọng.")

        # Ô nhập liệu
        if p := st.chat_input("Hỏi gì đó..."):
            st.session_state.msg.append({"role": "user", "content": p})
            with st.chat_message("user"): st.markdown(p)
            
            with st.chat_message("assistant"):
                # Dùng hàm an toàn thay vì gọi trực tiếp
                res = goi_gemini_an_toan(f"Dựa vào tài liệu:\n{st.session_state['noi_dung']}\nTrả lời: {p}")
                if res:
                    st.markdown(res.text)
                    st.session_state.msg.append({"role": "assistant", "content": res.text})

    # -----------------------------------------------------
    # TAB 2: TRẮC NGHIỆM (QUIZ)
    # -----------------------------------------------------
    with t2:
        c1, c2 = st.columns([1,3], vertical_alignment="bottom")
        sl = c1.number_input("Số câu", 1, 50, 5)
        
        if c2.button("🚀 Tạo Đề Trắc Nghiệm", use_container_width=True):
            with st.spinner("Đang tạo đề..."):
                try:
                    p = f"Tạo {sl} câu trắc nghiệm JSON list: [{{'question':'...','options':['A...'],'correct':'A','explain':'...'}}]"
                    res = goi_gemini_an_toan(f"{p}\nNội dung: {st.session_state['noi_dung']}")
                    if res:
                        st.session_state['quiz'] = json.loads(lay_json(res.text))
                except: st.error("Lỗi định dạng dữ liệu từ AI.")
        
        if 'quiz' in st.session_state:
            score = 0
            for i, q in enumerate(st.session_state['quiz']):
                st.divider()
                st.markdown(f"**{i+1}.** {q['question']}")
                ch = st.radio("Chọn đáp án:", q['options'], key=f"q{i}", index=None)
                if ch:
                    if ch[0] == q['correct'][0]:
                        st.success("✅ Chính xác!")
                        score+=1
                    else: 
                        st.error(f"❌ Sai rồi. Đáp án đúng: {q['correct']}")
                    with st.expander("🔍 Xem giải thích"): st.write(q['explain'])
            st.info(f"🏆 Điểm số của bạn: {score}/{len(st.session_state['quiz'])}")

    # -----------------------------------------------------
    # TAB 3: FLASHCARDS (ĐÃ NÂNG CẤP)
    # -----------------------------------------------------
    with t3:
        # [MỚI] CSS làm đẹp nút bấm cột 2
        st.markdown("""
            <style>
            div[data-testid="column"]:nth-of-type(2) .stButton > button {
                width: 100%;
                background: linear-gradient(to right, #4CAF50, #2E8B57);
                color: white;
                border-radius: 20px;
                height: 50px;
                border: none;
                font-weight: bold;
                box-shadow: 0 4px 6px rgba(0,0,0,0.2);
            }
            div[data-testid="column"]:nth-of-type(2) .stButton > button:hover {
                transform: scale(1.02);
                color: white;
            }
            </style>
        """, unsafe_allow_html=True)

        c1, c2 = st.columns([1,3], vertical_alignment="bottom")
        sl = c1.number_input("Số lượng thẻ", 1, 50, 5)
        
        if c2.button("🗂️ Tạo Flashcards Ngay", use_container_width=True):
            with st.spinner("Đang phân tích và tạo thẻ..."):
                try:
                    # [MỚI] Prompt chuẩn JSON với dấu ngoặc kép và {{ }}
                    p = f"Tạo {sl} cặp câu hỏi - đáp án ngắn gọn. Trả về JSON list thuần túy. Bắt buộc dùng dấu ngoặc kép cho Key và Value. Mẫu: [{{ \"q\": \"Câu hỏi?\", \"a\": \"Đáp án.\" }}]"
                    
                    res = goi_gemini_an_toan(f"{p}\n\nNội dung tài liệu:\n{st.session_state['noi_dung']}")
                    
                    if res:
                        st.session_state['fc'] = json.loads(lay_json(res.text))
                        st.success(f"✅ Đã tạo xong {len(st.session_state['fc'])} thẻ!")
                except Exception as e: 
                    st.error(f"Lỗi tạo thẻ: {e}")

        if 'fc' in st.session_state:
            st.write("---")
            for i, c in enumerate(st.session_state['fc']):
                with st.expander(f"🔹 Thẻ {i+1}: {c.get('q','?')}"): 
                    st.info(f"💡 {c.get('a','!')}")

    # -----------------------------------------------------
    # TAB 4: MINDMAP
    # -----------------------------------------------------
    with t4:
        st.subheader("Bản đồ kiến thức")
        if st.button("🎨 Vẽ Sơ Đồ Tư Duy"):
            with st.spinner("Đang vẽ sơ đồ..."):
                try:
                    p = """
                    Tóm tắt thành Sơ đồ tư duy (Mind Map).
                    Output format: chỉ mã Graphviz DOT (trong ```dot ... ```).
                    Dùng digraph G { rankdir="LR"; node [shape=box, style=filled, fillcolor="#E8F5E9", fontname="Arial"]; ... }
                    """
                    res = goi_gemini_an_toan(f"{p}\nNội dung: {st.session_state['noi_dung']}")
                    if res:
                        st.session_state['map'] = lay_dot_code(res.text)
                except: st.error("Lỗi vẽ hình.")
        
        if 'map' in st.session_state:
            try: st.graphviz_chart(st.session_state['map'])
            except: st.error("Mã hình lỗi, vui lòng thử lại.")

else:
    st.info("👈 Vui lòng tải tài liệu (PDF, Word, PPT...) ở thanh bên trái để bắt đầu!")
    st.warning("Lưu ý: Tải lại trang sẽ mất dữ liệu phiên làm việc hiện tại.")
